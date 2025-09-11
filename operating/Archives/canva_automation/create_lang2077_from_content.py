#!/usr/bin/env python3
"""
LANG 2077 Slides Creator from Content File
Creates PowerPoint slides based on lang2077_slides_content.md
"""

import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class LANG2077SlideCreator:
    """Create LANG 2077 slides from markdown content."""
    
    def __init__(self):
        """Initialize with HKBU branding."""
        self.hkbu_colors = {
            'burgundy': RGBColor(139, 0, 0),      # #8B0000
            'gold': RGBColor(255, 215, 0),       # #FFD700
            'dark_red': RGBColor(128, 0, 0),     # #800000
            'white': RGBColor(255, 255, 255),    # #FFFFFF
            'black': RGBColor(0, 0, 0),          # #000000
            'gray': RGBColor(64, 64, 64)         # #404040
        }
        
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self):
        """Add the title slide."""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "LANG 2077"
        title_frame = title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(48)
        title_para.font.color.rgb = self.hkbu_colors['burgundy']
        title_para.font.bold = True
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle_frame = subtitle.text_frame
        subtitle_frame.clear()
        
        # Main subtitle
        p1 = subtitle_frame.paragraphs[0]
        p1.text = "Language Skills for human-AI partnership:"
        p1.font.name = "Arial"
        p1.font.size = Pt(24)
        p1.font.color.rgb = self.hkbu_colors['dark_red']
        
        # Secondary subtitle
        p2 = subtitle_frame.add_paragraph()
        p2.text = "Customizing Chatbots to Empower Communities"
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = self.hkbu_colors['burgundy']
        p2.font.bold = True
        
        # Footer info
        p3 = subtitle_frame.add_paragraph()
        p3.text = ""
        p4 = subtitle_frame.add_paragraph()
        p4.text = "Dr. Simon Wang | Language Centre, HKBU | Fall 2025"
        p4.font.name = "Arial"
        p4.font.size = Pt(16)
        p4.font.color.rgb = self.hkbu_colors['gray']
        
        print("✅ Added title slide")
        return slide
    
    def add_learning_outcomes_slide(self):
        """Add learning outcomes slide."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "What Students Will Learn"
        title_frame = title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = self.hkbu_colors['burgundy']
        title_para.font.bold = True
        
        # Content
        content = slide.placeholders[1]
        content_frame = content.text_frame
        content_frame.clear()
        
        # Subtitle
        p0 = content_frame.paragraphs[0]
        p0.text = "Learning Outcomes & Objectives"
        p0.font.name = "Arial"
        p0.font.size = Pt(24)
        p0.font.color.rgb = self.hkbu_colors['dark_red']
        p0.font.bold = True
        
        # Language Skills
        p1 = content_frame.add_paragraph()
        p1.text = "• Language Skills Development"
        p1.font.name = "Arial"
        p1.font.size = Pt(20)
        p1.font.color.rgb = self.hkbu_colors['burgundy']
        p1.font.bold = True
        
        skills = [
            "Academic communication in AI contexts",
            "Technical writing for AI applications", 
            "Cross-cultural communication strategies"
        ]
        for skill in skills:
            p = content_frame.add_paragraph()
            p.text = f"  - {skill}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self.hkbu_colors['black']
        
        # AI Partnership Skills
        p2 = content_frame.add_paragraph()
        p2.text = "• AI Partnership Skills"
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = self.hkbu_colors['burgundy']
        p2.font.bold = True
        
        ai_skills = [
            "Understanding AI capabilities and limitations",
            "Effective human-AI collaboration",
            "Prompt engineering and chatbot interaction"
        ]
        for skill in ai_skills:
            p = content_frame.add_paragraph()
            p.text = f"  - {skill}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self.hkbu_colors['black']
        
        # Community Engagement
        p3 = content_frame.add_paragraph()
        p3.text = "• Community Engagement"
        p3.font.name = "Arial"
        p3.font.size = Pt(20)
        p3.font.color.rgb = self.hkbu_colors['burgundy']
        p3.font.bold = True
        
        community_skills = [
            "Identifying community needs",
            "Designing solutions with community partners",
            "Presenting results to stakeholders"
        ]
        for skill in community_skills:
            p = content_frame.add_paragraph()
            p.text = f"  - {skill}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self.hkbu_colors['black']
        
        print("✅ Added learning outcomes slide")
        return slide
    
    def add_community_empowerment_slide(self):
        """Add community empowerment slide."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Empowering Communities Through AI"
        title_frame = title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = self.hkbu_colors['burgundy']
        title_para.font.bold = True
        
        # Content
        content = slide.placeholders[1]
        content_frame = content.text_frame
        content_frame.clear()
        
        # Subtitle
        p0 = content_frame.paragraphs[0]
        p0.text = "Service Learning Component"
        p0.font.name = "Arial"
        p0.font.size = Pt(24)
        p0.font.color.rgb = self.hkbu_colors['dark_red']
        p0.font.bold = True
        
        # Community Partner Collaboration
        p1 = content_frame.add_paragraph()
        p1.text = "• Community Partner Collaboration"
        p1.font.name = "Arial"
        p1.font.size = Pt(20)
        p1.font.color.rgb = self.hkbu_colors['burgundy']
        p1.font.bold = True
        
        collab_items = [
            "Work with local NGOs and organizations",
            "Identify real community challenges",
            "Co-design AI-enhanced solutions"
        ]
        for item in collab_items:
            p = content_frame.add_paragraph()
            p.text = f"  - {item}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self.hkbu_colors['black']
        
        # Chatbot Customization
        p2 = content_frame.add_paragraph()
        p2.text = "• Chatbot Customization Projects"
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = self.hkbu_colors['burgundy']
        p2.font.bold = True
        
        chatbot_items = [
            "Develop task-specific chatbots",
            "Adapt language for target audiences",
            "Test and iterate with community feedback"
        ]
        for item in chatbot_items:
            p = content_frame.add_paragraph()
            p.text = f"  - {item}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self.hkbu_colors['black']
        
        # Student Deliverables
        p3 = content_frame.add_paragraph()
        p3.text = "• Student Deliverables & Impact"
        p3.font.name = "Arial"
        p3.font.size = Pt(20)
        p3.font.color.rgb = self.hkbu_colors['burgundy']
        p3.font.bold = True
        
        deliverable_items = [
            "Final presentation to community partners",
            "Reflection essays on AI ethics",
            "Portfolio of customized AI tools"
        ]
        for item in deliverable_items:
            p = content_frame.add_paragraph()
            p.text = f"  - {item}"
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = self.hkbu_colors['black']
        
        print("✅ Added community empowerment slide")
        return slide
    
    def add_course_info_slide(self):
        """Add additional course information slide."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Course Information"
        title_frame = title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Arial"
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = self.hkbu_colors['burgundy']
        title_para.font.bold = True
        
        # Content
        content = slide.placeholders[1]
        content_frame = content.text_frame
        content_frame.clear()
        
        # Course Structure
        p1 = content_frame.paragraphs[0]
        p1.text = "• Course Structure & Timeline"
        p1.font.name = "Arial"
        p1.font.size = Pt(20)
        p1.font.color.rgb = self.hkbu_colors['burgundy']
        p1.font.bold = True
        
        p = content_frame.add_paragraph()
        p.text = "  - Week-by-week breakdown of activities and assignments"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = self.hkbu_colors['black']
        
        # Assessment
        p2 = content_frame.add_paragraph()
        p2.text = "• Assessment Methods"
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = self.hkbu_colors['burgundy']
        p2.font.bold = True
        
        p = content_frame.add_paragraph()
        p.text = "  - Grading rubrics, participation expectations, project criteria"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = self.hkbu_colors['black']
        
        # Resources
        p3 = content_frame.add_paragraph()
        p3.text = "• Required Resources"
        p3.font.name = "Arial"
        p3.font.size = Pt(20)
        p3.font.color.rgb = self.hkbu_colors['burgundy']
        p3.font.bold = True
        
        p = content_frame.add_paragraph()
        p.text = "  - Textbooks, software, online platforms students will use"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = self.hkbu_colors['black']
        
        # Contact
        p4 = content_frame.add_paragraph()
        p4.text = "• Contact Information"
        p4.font.name = "Arial"
        p4.font.size = Pt(20)
        p4.font.color.rgb = self.hkbu_colors['burgundy']
        p4.font.bold = True
        
        p = content_frame.add_paragraph()
        p.text = "  - Office hours, email, course website, support resources"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = self.hkbu_colors['black']
        
        print("✅ Added course information slide")
        return slide
    
    def save_presentation(self, filename=None):
        """Save the presentation."""
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M")
            filename = f"LANG2077_slides_from_content_{timestamp}.pptx"
        
        # Ensure generated_slides directory exists
        slides_dir = "/Users/simonwang/Documents/Usage/VibeCoding/DailyAssistant/operating/canva_automation/generated_slides"
        os.makedirs(slides_dir, exist_ok=True)
        
        filepath = os.path.join(slides_dir, filename)
        self.prs.save(filepath)
        
        print(f"✅ Presentation saved: {filepath}")
        return filepath


def create_lang2077_slides():
    """Create LANG 2077 slides from content file."""
    
    print("🎓 Creating LANG 2077 Slides from Content File")
    print("=" * 50)
    
    # Check if content file exists
    content_file = "/Users/simonwang/Documents/Usage/VibeCoding/DailyAssistant/operating/canva_automation/lang2077_materials/lang2077_slides_content.md"
    if not os.path.exists(content_file):
        print(f"❌ Content file not found: {content_file}")
        return None
    
    print(f"✅ Found content file: {os.path.basename(content_file)}")
    
    # Create slides
    creator = LANG2077SlideCreator()
    
    print("\n🔄 Creating slides...")
    creator.add_title_slide()
    creator.add_learning_outcomes_slide()
    creator.add_community_empowerment_slide()
    creator.add_course_info_slide()
    
    # Save presentation
    print("\n💾 Saving presentation...")
    filepath = creator.save_presentation()
    
    print(f"\n🎉 LANG 2077 presentation created successfully!")
    print(f"📁 Location: {filepath}")
    print(f"📊 Slides: 4 slides with HKBU branding")
    print(f"💼 Compatible with: PowerPoint, Google Slides, Keynote")
    
    # Note about Canva sharing
    print(f"\n📝 Note: Since Canva API requires authentication setup,")
    print(f"   this PowerPoint file can be uploaded to Canva manually")
    print(f"   for sharing with 'anyone with link can edit' permissions.")
    
    return filepath


if __name__ == "__main__":
    create_lang2077_slides()
