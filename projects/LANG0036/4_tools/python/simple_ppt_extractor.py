"""
Simple PowerPoint Extractor
Quick script to extract text and images from the specific PowerPoint file
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from pathlib import Path


def extract_ppt_simple():
    """Simple extraction function"""
    
    # File path
    ppt_file = r"C:\Users\simonwang\Google Drive Streaming\Other computers\My Mac\VibeCoding\DailyAssistant\projects\GCAP3226\course_materials\week4\GCAP3226_week4_Simulation_20250920.pptx"
    
    # Output directory
    output_dir = Path(ppt_file).parent / "extracted_ppt_content"
    output_dir.mkdir(exist_ok=True)
    
    print(f"Loading PowerPoint: {ppt_file}")
    
    try:
        # Load presentation
        prs = Presentation(ppt_file)
        
        # Extract text
        all_text = []
        image_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\nSlide {slide_num}:")
            slide_text = f"=== SLIDE {slide_num} ===\n"
            
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text += f"{shape.text.strip()}\n"
                    print(f"  Text found: {shape.text[:50]}...")
                
                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Save image
                        image_filename = f"slide_{slide_num:02d}_image_{image_count:02d}.png"
                        image_path = output_dir / image_filename
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        
                        print(f"  Image saved: {image_filename}")
                        image_count += 1
                        
                    except Exception as e:
                        print(f"  Error saving image: {e}")
            
            all_text.append(slide_text + "\n" + "-"*50 + "\n")
        
        # Save all text
        text_file = output_dir / "extracted_text.txt"
        with open(text_file, 'w', encoding='utf-8') as f:
            f.write("".join(all_text))
        
        print(f"\n✅ Extraction complete!")
        print(f"📄 Text saved to: {text_file}")
        print(f"🖼️  {image_count} images saved to: {output_dir}")
        print(f"📁 Output directory: {output_dir}")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        print("Make sure you have installed: pip install python-pptx")


if __name__ == "__main__":
    extract_ppt_simple()