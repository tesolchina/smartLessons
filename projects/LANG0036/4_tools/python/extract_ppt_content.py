"""
PowerPoint Content Extractor
Extracts text and images from PowerPoint presentations

Requirements: pip install python-pptx pillow
"""

import os
import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import json


def extract_text_from_slide(slide):
    """Extract all text content from a slide"""
    text_content = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            text_content.append(shape.text.strip())
        
        # Handle tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_text = []
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text:
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text.append(' | '.join(row_text))
            if table_text:
                text_content.append('\n'.join(table_text))
    
    return text_content


def extract_images_from_slide(slide, slide_num, output_dir):
    """Extract all images from a slide"""
    images_info = []
    image_count = 0
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                
                # Determine file extension
                ext = 'png'  # default
                if image.content_type == 'image/jpeg':
                    ext = 'jpg'
                elif image.content_type == 'image/png':
                    ext = 'png'
                elif image.content_type == 'image/gif':
                    ext = 'gif'
                
                # Create filename
                filename = f"slide_{slide_num:02d}_image_{image_count:02d}.{ext}"
                filepath = output_dir / filename
                
                # Save image
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
                
                # Try to get image dimensions
                try:
                    img = Image.open(io.BytesIO(image_bytes))
                    width, height = img.size
                except:
                    width, height = None, None
                
                images_info.append({
                    'filename': filename,
                    'filepath': str(filepath),
                    'content_type': image.content_type,
                    'size_bytes': len(image_bytes),
                    'dimensions': {'width': width, 'height': height} if width and height else None
                })
                
                image_count += 1
                print(f"  Extracted image: {filename}")
                
            except Exception as e:
                print(f"  Error extracting image from slide {slide_num}: {e}")
    
    return images_info


def extract_ppt_content(ppt_path, output_dir=None):
    """
    Extract all text and images from a PowerPoint presentation
    
    Args:
        ppt_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save extracted content (optional)
    
    Returns:
        dict: Extracted content information
    """
    ppt_path = Path(ppt_path)
    
    if not ppt_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")
    
    # Set up output directory
    if output_dir is None:
        output_dir = ppt_path.parent / f"{ppt_path.stem}_extracted"
    else:
        output_dir = Path(output_dir)
    
    # Create directories
    images_dir = output_dir / "images"
    text_dir = output_dir / "text"
    
    images_dir.mkdir(parents=True, exist_ok=True)
    text_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"Extracting content from: {ppt_path}")
    print(f"Output directory: {output_dir}")
    
    try:
        # Load presentation
        prs = Presentation(ppt_path)
        
        # Initialize results
        extraction_results = {
            'source_file': str(ppt_path),
            'output_directory': str(output_dir),
            'total_slides': len(prs.slides),
            'slides': [],
            'summary': {
                'total_text_blocks': 0,
                'total_images': 0,
                'slides_with_text': 0,
                'slides_with_images': 0
            }
        }
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\nProcessing slide {slide_num}...")
            
            # Extract text
            slide_texts = extract_text_from_slide(slide)
            
            # Extract images
            slide_images = extract_images_from_slide(slide, slide_num, images_dir)
            
            # Save slide text to individual file
            slide_text_file = text_dir / f"slide_{slide_num:02d}.txt"
            with open(slide_text_file, 'w', encoding='utf-8') as f:
                f.write(f"=== SLIDE {slide_num} ===\n\n")
                for i, text in enumerate(slide_texts, 1):
                    f.write(f"Text Block {i}:\n{text}\n\n")
            
            # Store slide information
            slide_info = {
                'slide_number': slide_num,
                'text_blocks': slide_texts,
                'text_file': str(slide_text_file),
                'images': slide_images,
                'text_count': len(slide_texts),
                'image_count': len(slide_images)
            }
            
            extraction_results['slides'].append(slide_info)
            
            # Update summary
            extraction_results['summary']['total_text_blocks'] += len(slide_texts)
            extraction_results['summary']['total_images'] += len(slide_images)
            if slide_texts:
                extraction_results['summary']['slides_with_text'] += 1
            if slide_images:
                extraction_results['summary']['slides_with_images'] += 1
            
            print(f"  Text blocks: {len(slide_texts)}")
            print(f"  Images: {len(slide_images)}")
        
        # Save combined text file
        combined_text_file = text_dir / "all_slides_combined.txt"
        with open(combined_text_file, 'w', encoding='utf-8') as f:
            f.write(f"EXTRACTED TEXT FROM: {ppt_path.name}\n")
            f.write(f"{'='*50}\n\n")
            
            for slide_info in extraction_results['slides']:
                f.write(f"=== SLIDE {slide_info['slide_number']} ===\n\n")
                for i, text in enumerate(slide_info['text_blocks'], 1):
                    f.write(f"Text Block {i}:\n{text}\n\n")
                f.write(f"{'-'*30}\n\n")
        
        # Save extraction metadata as JSON
        metadata_file = output_dir / "extraction_metadata.json"
        with open(metadata_file, 'w', encoding='utf-8') as f:
            json.dump(extraction_results, f, indent=2, ensure_ascii=False)
        
        print(f"\n✅ Extraction completed successfully!")
        print(f"📄 Combined text saved to: {combined_text_file}")
        print(f"🖼️  Images saved to: {images_dir}")
        print(f"📊 Metadata saved to: {metadata_file}")
        print(f"\nSummary:")
        print(f"  Total slides: {extraction_results['summary']['total_text_blocks']}")
        print(f"  Text blocks: {extraction_results['summary']['total_text_blocks']}")
        print(f"  Images: {extraction_results['summary']['total_images']}")
        
        return extraction_results
        
    except Exception as e:
        print(f"❌ Error processing PowerPoint file: {e}")
        raise


def main():
    """Main function to run the extraction"""
    
    # Configuration
    ppt_file = r"C:\Users\simonwang\Google Drive Streaming\Other computers\My Mac\VibeCoding\DailyAssistant\projects\GCAP3226\course_materials\week4\GCAP3226_week4_Simulation_20250920.pptx"
    
    # Optional: specify custom output directory
    # output_directory = r"C:\path\to\custom\output"
    output_directory = None  # Will create a folder next to the PPT file
    
    try:
        # Extract content
        results = extract_ppt_content(ppt_file, output_directory)
        
        print(f"\n🎉 All done! Check the output directory for extracted content.")
        
    except FileNotFoundError as e:
        print(f"❌ File error: {e}")
    except Exception as e:
        print(f"❌ Unexpected error: {e}")


if __name__ == "__main__":
    main()