#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for HKBU Flipped Classroom Workshop
Creates a comprehensive presentation based on the workshop planning document
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os

def create_workshop_presentation():
    """Create the main workshop presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions to widescreen
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Define color scheme based on your workshop theme
    colors = {
        'primary_orange': RGBColor(255, 107, 53),     # #FF6B35
        'secondary_orange': RGBColor(255, 179, 128),   # #FFB380
        'primary_purple': RGBColor(139, 92, 246),      # #8B5CF6
        'secondary_purple': RGBColor(167, 139, 250),   # #A78BFA
        'warm_yellow': RGBColor(254, 243, 199),        # #FEF3C7
        'warm_peach': RGBColor(255, 238, 230),         # #FFEEE6
        'dark_text': RGBColor(31, 41, 55),             # #1F2937
        'medium_text': RGBColor(75, 85, 99),           # #4B5563
        'white': RGBColor(255, 255, 255)
    }
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_slide(slide1, colors)
    
    # Slide 2: Academic Foundation
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_academic_foundation_slide(slide2, colors)
    
    # Slide 3: The Challenge
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_challenge_slide(slide3, colors)
    
    # Slide 4: AI Solution
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_ai_solution_slide(slide4, colors)
    
    # Slide 5: Workshop Journey
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_workshop_journey_slide(slide5, colors)
    
    # Slide 6: Call to Action
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_call_to_action_slide(slide6, colors)
    
    return prs

def add_gradient_background(slide, color1, color2):
    """Add gradient background to slide"""
    # Note: python-pptx doesn't directly support gradients, so we'll use solid colors
    # You can manually add gradients in PowerPoint after generation
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1

def add_title_slide(slide, colors):
    """Create the title slide"""
    
    # Add gradient background effect with shapes
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9)
    )
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = colors['warm_peach']
    bg_shape.line.fill.background()
    
    # Add decorative shapes
    shape1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(13), Inches(1), Inches(4), Inches(4))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = colors['primary_purple']
    shape1.fill.transparency = 0.8
    shape1.line.fill.background()
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(6), Inches(3), Inches(3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = colors['primary_orange']
    shape2.fill.transparency = 0.8
    shape2.line.fill.background()
    
    # Badge
    badge = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(3), Inches(0.8))
    badge_frame = badge.text_frame
    badge_frame.text = "Let's Flip It! 🚀"
    badge_p = badge_frame.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    badge_font = badge_p.runs[0].font
    badge_font.size = Pt(16)
    badge_font.bold = True
    badge_font.color.rgb = colors['white']
    
    # Badge background
    badge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(1.4), Inches(3.4), Inches(1))
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = colors['primary_purple']
    badge_bg.line.fill.background()
    # Move badge background behind text
    badge_bg_element = badge_bg.element
    slide.shapes._spTree.remove(badge_bg_element)
    slide.shapes._spTree.insert(2, badge_bg_element)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Engaging Students in Flipped Classroom with AI-Powered Video and Avatars"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_font = title_p.runs[0].font
    title_font.size = Pt(36)
    title_font.bold = True
    title_font.color.rgb = colors['dark_text']
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Innovative AI-assisted tools for enhanced learning experiences"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_font = subtitle_p.runs[0].font
    subtitle_font.size = Pt(20)
    subtitle_font.color.rgb = colors['medium_text']
    
    # Workshop details
    details_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(14), Inches(1))
    details_frame = details_box.text_frame
    details_frame.text = "September 26, 2025 | 2:30 PM | ZOOM"
    details_p = details_frame.paragraphs[0]
    details_p.alignment = PP_ALIGN.CENTER
    details_font = details_p.runs[0].font
    details_font.size = Pt(18)
    details_font.bold = True
    details_font.color.rgb = colors['primary_purple']
    
    # Speaker info background
    speaker_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(6.8), Inches(8), Inches(1.5))
    speaker_bg.fill.solid()
    speaker_bg.fill.fore_color.rgb = colors['warm_yellow']
    speaker_bg.line.fill.background()
    
    # Speaker info
    speaker_box = slide.shapes.add_textbox(Inches(4.2), Inches(7), Inches(7.6), Inches(1.1))
    speaker_frame = speaker_box.text_frame
    speaker_frame.text = "Dr Simon Wang\nLecturer in English & Innovation Officer\nLanguage Centre, HKBU"
    speaker_p = speaker_frame.paragraphs[0]
    speaker_p.alignment = PP_ALIGN.CENTER
    speaker_font = speaker_p.runs[0].font
    speaker_font.size = Pt(16)
    speaker_font.bold = True
    speaker_font.color.rgb = colors['primary_purple']

def add_academic_foundation_slide(slide, colors):
    """Create the academic foundation slide"""
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = colors['white']
    bg_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "📚 Built on Solid Academic Research"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_font = title_p.runs[0].font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = colors['dark_text']
    
    # Statistics grid
    stats = [
        ("2025", "Latest AI Integration\nStudies"),
        ("2017", "HKBU Flipped\nClassroom Journey"),
        ("8+", "Years of\nProven Results")
    ]
    
    for i, (number, description) in enumerate(stats):
        x_pos = Inches(1.5 + i * 4.5)
        
        # Stat card background
        stat_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(4), Inches(2))
        stat_bg.fill.solid()
        stat_bg.fill.fore_color.rgb = colors['warm_peach']
        stat_bg.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(x_pos, Inches(2.2), Inches(4), Inches(0.8))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.alignment = PP_ALIGN.CENTER
        num_font = num_p.runs[0].font
        num_font.size = Pt(28)
        num_font.bold = True
        num_font.color.rgb = colors['primary_orange']
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos, Inches(3), Inches(4), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_p = desc_frame.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER
        desc_font = desc_p.runs[0].font
        desc_font.size = Pt(14)
        desc_font.color.rgb = colors['medium_text']
    
    # Research foundation box
    research_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.8), Inches(14), Inches(3.5))
    research_bg.fill.solid()
    research_bg.fill.fore_color.rgb = colors['warm_yellow']
    research_bg.line.fill.background()
    
    # Research content
    research_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(13), Inches(3))
    research_frame = research_box.text_frame
    research_frame.text = """🏛️ Research Foundation

• "Integrating artificial intelligence in the flipped classroom" (2025)
• "AI-assisted assessment in higher education" (2024)  
• "INTRODUCING AI IN FLIPPED CLASSROOM" (2024)
• HKBU's SCMP Letter - Pioneering flexible learning (2017)

Our approach is backed by peer-reviewed research and proven results at HKBU."""
    
    # Format research text
    for i, paragraph in enumerate(research_frame.paragraphs):
        if i == 0:  # Title
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors['primary_purple']
        else:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = colors['dark_text']

def add_challenge_slide(slide, colors):
    """Create the challenge slide"""
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = colors['white']
    bg_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "⚡ The Traditional Challenge"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_font = title_p.runs[0].font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = colors['dark_text']
    
    # Challenge statistics
    stats = [
        ("4-6", "Hours per\nvideo"),
        ("75%", "Technical\nbarriers"),
        ("320+", "Hours\nannually")
    ]
    
    for i, (number, description) in enumerate(stats):
        x_pos = Inches(1.5 + i * 4.5)
        
        # Stat card background (red-ish for problems)
        stat_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(4), Inches(2))
        stat_bg.fill.solid()
        stat_bg.fill.fore_color.rgb = RGBColor(255, 200, 200)  # Light red
        stat_bg.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(x_pos, Inches(2.2), Inches(4), Inches(0.8))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.alignment = PP_ALIGN.CENTER
        num_font = num_p.runs[0].font
        num_font.size = Pt(28)
        num_font.bold = True
        num_font.color.rgb = RGBColor(220, 50, 50)  # Dark red
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos, Inches(3), Inches(4), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_p = desc_frame.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER
        desc_font = desc_p.runs[0].font
        desc_font.size = Pt(14)
        desc_font.color.rgb = colors['medium_text']
    
    # Pain points box
    pain_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.8), Inches(14), Inches(3.5))
    pain_bg.fill.solid()
    pain_bg.fill.fore_color.rgb = RGBColor(255, 240, 240)  # Very light red
    pain_bg.line.fill.background()
    
    # Pain points content
    pain_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(13), Inches(3))
    pain_frame = pain_box.text_frame
    pain_frame.text = """🔥 The Pain Points

• Video editing complexity requiring specialized skills
• Time-intensive production taking away from teaching and research
• Technical skill requirements creating barriers to adoption
• Limited scalability for multiple courses and content updates"""
    
    # Format pain points text
    for i, paragraph in enumerate(pain_frame.paragraphs):
        if i == 0:  # Title
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(180, 50, 50)
        else:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = colors['dark_text']

def add_ai_solution_slide(slide, colors):
    """Create the AI solution slide"""
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = colors['white']
    bg_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 The AI-Powered Solution"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_font = title_p.runs[0].font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = colors['dark_text']
    
    # Solution statistics
    stats = [
        ("5-10", "Minutes per\nvideo"),
        ("95%", "Time\nsavings"),
        ("24/7", "Avatar\nsupport")
    ]
    
    for i, (number, description) in enumerate(stats):
        x_pos = Inches(1.5 + i * 4.5)
        
        # Stat card background (green for solutions)
        stat_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(4), Inches(2))
        stat_bg.fill.solid()
        stat_bg.fill.fore_color.rgb = RGBColor(200, 255, 200)  # Light green
        stat_bg.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(x_pos, Inches(2.2), Inches(4), Inches(0.8))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.alignment = PP_ALIGN.CENTER
        num_font = num_p.runs[0].font
        num_font.size = Pt(28)
        num_font.bold = True
        num_font.color.rgb = RGBColor(50, 150, 50)  # Dark green
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos, Inches(3), Inches(4), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_p = desc_frame.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER
        desc_font = desc_p.runs[0].font
        desc_font.size = Pt(14)
        desc_font.color.rgb = colors['medium_text']
    
    # AI superpowers box
    ai_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.8), Inches(14), Inches(3.5))
    ai_bg.fill.solid()
    ai_bg.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Very light green
    ai_bg.line.fill.background()
    
    # AI superpowers content
    ai_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(13), Inches(3))
    ai_frame = ai_box.text_frame
    ai_frame.text = """✨ AI Superpowers

🎥 Instant video generation with professional quality
🤖 Intelligent avatars providing 24/7 student support
💻 Vibe coding tools for easy customization
📊 Automatic analytics and performance tracking"""
    
    # Format AI superpowers text
    for i, paragraph in enumerate(ai_frame.paragraphs):
        if i == 0:  # Title
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors['primary_purple']
        else:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = colors['dark_text']

def add_workshop_journey_slide(slide, colors):
    """Create the workshop journey slide"""
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = colors['white']
    bg_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "🎓 Your Learning Journey Today"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_font = title_p.runs[0].font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = colors['dark_text']
    
    # Learning modules
    modules = [
        ("🎬", "AI Video\nCreation"),
        ("🤖", "Streaming\nAvatars"),
        ("💻", "Vibe\nCoding"),
        ("⏰", "Time\nOptimization"),
        ("🚀", "Implementation"),
        ("📚", "Resources &\nNext Steps")
    ]
    
    # Create 2x3 grid
    for i, (icon, title) in enumerate(modules):
        row = i // 3
        col = i % 3
        x_pos = Inches(1.5 + col * 4.5)
        y_pos = Inches(2.2 + row * 2.2)
        
        # Module card background
        module_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, Inches(4), Inches(1.8))
        module_bg.fill.solid()
        module_bg.fill.fore_color.rgb = colors['warm_peach']
        module_bg.line.fill.background()
        
        # Icon
        icon_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.1), Inches(4), Inches(0.7))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_p = icon_frame.paragraphs[0]
        icon_p.alignment = PP_ALIGN.CENTER
        icon_font = icon_p.runs[0].font
        icon_font.size = Pt(24)
        
        # Title
        title_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.8), Inches(4), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_font = title_p.runs[0].font
        title_font.size = Pt(14)
        title_font.bold = True
        title_font.color.rgb = colors['primary_purple']
    
    # Objectives box
    obj_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6.8), Inches(14), Inches(1.8))
    obj_bg.fill.solid()
    obj_bg.fill.fore_color.rgb = colors['warm_yellow']
    obj_bg.line.fill.background()
    
    # Objectives content
    obj_box = slide.shapes.add_textbox(Inches(1.5), Inches(7), Inches(13), Inches(1.4))
    obj_frame = obj_box.text_frame
    obj_frame.text = """🎯 By the end of today, you'll be able to:
✅ Create your first AI-generated educational video  ✅ Deploy a 24/7 AI teaching assistant
✅ Experience AI-assisted coding (no programming background needed)  ✅ Calculate your time savings and ROI"""
    
    obj_p = obj_frame.paragraphs[0]
    obj_font = obj_p.runs[0].font
    obj_font.size = Pt(16)
    obj_font.bold = True
    obj_font.color.rgb = colors['primary_purple']

def add_call_to_action_slide(slide, colors):
    """Create the call to action slide"""
    
    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = colors['white']
    bg_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 Join the HKBU Innovation Community"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_font = title_p.runs[0].font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = colors['dark_text']
    
    # Quote box
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(12), Inches(1.5))
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = colors['primary_purple']
    quote_bg.line.fill.background()
    
    quote_box = slide.shapes.add_textbox(Inches(2.2), Inches(2.2), Inches(11.6), Inches(1.1))
    quote_frame = quote_box.text_frame
    quote_frame.text = '"AI doesn\'t replace teachers - it amplifies their impact."\nStop building tools, start teaching and researching!'
    quote_p = quote_frame.paragraphs[0]
    quote_p.alignment = PP_ALIGN.CENTER
    quote_font = quote_p.runs[0].font
    quote_font.size = Pt(20)
    quote_font.bold = True
    quote_font.color.rgb = colors['white']
    
    # Impact statistics
    stats = [
        ("95%", "Time\nSavings"),
        ("24/7", "Student\nSupport"),
        ("∞", "Creative\nPossibilities")
    ]
    
    for i, (number, description) in enumerate(stats):
        x_pos = Inches(1.5 + i * 4.5)
        
        # Stat card background
        stat_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(4), Inches(4), Inches(1.5))
        stat_bg.fill.solid()
        stat_bg.fill.fore_color.rgb = colors['warm_peach']
        stat_bg.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(x_pos, Inches(4.1), Inches(4), Inches(0.6))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_p = num_frame.paragraphs[0]
        num_p.alignment = PP_ALIGN.CENTER
        num_font = num_p.runs[0].font
        num_font.size = Pt(24)
        num_font.bold = True
        num_font.color.rgb = colors['primary_orange']
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos, Inches(4.7), Inches(4), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_p = desc_frame.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER
        desc_font = desc_p.runs[0].font
        desc_font.size = Pt(12)
        desc_font.color.rgb = colors['medium_text']
    
    # Call to action box
    cta_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6), Inches(12), Inches(2.5))
    cta_bg.fill.solid()
    cta_bg.fill.fore_color.rgb = colors['warm_yellow']
    cta_bg.line.fill.background()
    
    cta_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.2), Inches(11), Inches(2.1))
    cta_frame = cta_box.text_frame
    cta_frame.text = """🎯 Ready to Transform Your Teaching?

✅ Sign up for the pilot program
✅ Schedule your one-on-one consultation  
✅ Join our innovation community

Let's explore the interactive modules together!"""
    
    # Format CTA text
    for i, paragraph in enumerate(cta_frame.paragraphs):
        if i == 0:  # Title
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors['primary_purple']
        else:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = colors['dark_text']

def main():
    """Main function to generate the presentation"""
    
    print("🎬 Generating HKBU Flipped Classroom Workshop Presentation...")
    
    # Create the presentation
    prs = create_workshop_presentation()
    
    # Save the presentation
    output_dir = "presentations"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    output_file = os.path.join(output_dir, "HKBU_Flipped_Classroom_Workshop_Overview.pptx")
    prs.save(output_file)
    
    print(f"✅ Presentation saved to: {output_file}")
    print(f"📊 Total slides created: {len(prs.slides)}")
    print()
    print("📝 Presentation Contents:")
    slide_titles = [
        "1. Title Slide - Workshop Introduction",
        "2. Academic Foundation - Research backing",
        "3. The Challenge - Traditional problems", 
        "4. AI Solution - Technology benefits",
        "5. Workshop Journey - Learning modules",
        "6. Call to Action - Next steps"
    ]
    
    for title in slide_titles:
        print(f"   {title}")
    
    print()
    print("🎯 Next steps:")
    print("   • Open the presentation in PowerPoint or Google Slides")
    print("   • Customize colors, fonts, and content as needed")
    print("   • Add animations and transitions")
    print("   • Embed in your overview.html component")

if __name__ == "__main__":
    main()
