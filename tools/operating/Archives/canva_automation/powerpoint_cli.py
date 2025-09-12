#!/usr/bin/env python3
"""
PowerPoint CLI Generator - Creates .pptx files programmatically
No external APIs required - generates native PowerPoint files
Created: September 6, 2025
"""

import os
import sys
import argparse
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
except ImportError:
    print("❌ Missing python-pptx dependency")
    print("📦 Install with: pip3 install python-pptx")
    sys.exit(1)

class PowerPointCLI:
    """Generate PowerPoint presentations programmatically"""
    
    def __init__(self):
        self.output_dir = "/Users/simonwang/Documents/Usage/VibeCoding/DailyAssistant/operating/canva_automation/generated_slides"
        Path(self.output_dir).mkdir(exist_ok=True)
        
        # HKBU Brand Colors (RGB)
        self.colors = {
            "burgundy": RGBColor(128, 0, 32),    # #800020
            "gold": RGBColor(255, 215, 0),       # #FFD700  
            "dark_gray": RGBColor(51, 51, 51),   # #333333
            "light_gray": RGBColor(245, 245, 245) # #F5F5F5
        }
    
    def create_lang2077_presentation(self, collaborators=None):
        """Create LANG 2077 PowerPoint presentation"""
        
        print("🎓 Creating LANG 2077 PowerPoint presentation...")
        
        # Create presentation object
        prs = Presentation()
        
        # Slide 1: Title Slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
        
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        
        title1.text = "LANG 2077"
        title1.text_frame.paragraphs[0].font.size = Pt(54)
        title1.text_frame.paragraphs[0].font.color.rgb = self.colors["burgundy"]
        title1.text_frame.paragraphs[0].font.bold = True
        
        subtitle1.text = """Language Skills for human-AI partnership:
Customizing Chatbots to Empower Communities

Dr. Simon Wang | Language Centre, HKBU | Fall 2025"""
        subtitle1.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle1.text_frame.paragraphs[1].font.size = Pt(24)
        subtitle1.text_frame.paragraphs[2].font.size = Pt(18)
        subtitle1.text_frame.paragraphs[2].font.color.rgb = self.colors["gold"]
        
        # Slide 2: Learning Outcomes
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        title2 = slide2.shapes.title
        title2.text = "What Students Will Learn"
        title2.text_frame.paragraphs[0].font.size = Pt(44)
        title2.text_frame.paragraphs[0].font.color.rgb = self.colors["burgundy"]
        
        content2 = slide2.placeholders[1]
        content2.text = """Language Skills Development
• Academic communication in AI contexts
• Technical writing for AI applications
• Cross-cultural communication strategies

AI Partnership Skills
• Understanding AI capabilities and limitations
• Effective human-AI collaboration  
• Prompt engineering and chatbot interaction

Community Engagement
• Identifying community needs
• Designing solutions with community partners
• Presenting results to stakeholders"""
        
        # Format content with colors
        for i, paragraph in enumerate(content2.text_frame.paragraphs):
            if i in [0, 4, 8]:  # Main headers
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = self.colors["burgundy"]
                paragraph.font.bold = True
            else:  # Bullet points
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = self.colors["dark_gray"]
        
        # Slide 3: Community Impact
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        title3 = slide3.shapes.title
        title3.text = "Empowering Communities Through AI"
        title3.text_frame.paragraphs[0].font.size = Pt(40)
        title3.text_frame.paragraphs[0].font.color.rgb = self.colors["burgundy"]
        
        content3 = slide3.placeholders[1]
        content3.text = """Community Partner Collaboration
• Work with local NGOs and organizations
• Identify real community challenges
• Co-design AI-enhanced solutions

Chatbot Customization Projects
• Develop task-specific chatbots
• Adapt language for target audiences
• Test and iterate with community feedback

Student Deliverables & Impact
• Final presentation to community partners
• Reflection essays on AI ethics
• Portfolio of customized AI tools"""
        
        # Format content
        for i, paragraph in enumerate(content3.text_frame.paragraphs):
            if i in [0, 4, 8]:  # Main headers
                paragraph.font.size = Pt(26)
                paragraph.font.color.rgb = self.colors["gold"]
                paragraph.font.bold = True
            else:  # Bullet points
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = self.colors["dark_gray"]
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        filename = f"LANG2077_presentation_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        
        prs.save(filepath)
        
        print(f"✅ PowerPoint presentation created!")
        print(f"📁 File: {filepath}")
        print(f"💼 Compatible with: PowerPoint, Google Slides, Keynote")
        
        # Create collaboration instructions if collaborators specified
        if collaborators:
            self._create_pptx_collaboration_guide(filepath, collaborators)
        
        return filepath
    
    def create_custom_presentation(self, title, slides_data, collaborators=None):
        """Create custom PowerPoint presentation"""
        
        print(f"🎨 Creating custom presentation: {title}")
        
        prs = Presentation()
        
        for i, slide_data in enumerate(slides_data):
            if slide_data.get("type") == "title":
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = slide_data["title"]
                if slide_data.get("subtitle"):
                    slide.placeholders[1].text = slide_data["subtitle"]
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data["title"]
                if slide_data.get("content"):
                    content_text = "\n".join(slide_data["content"]) if isinstance(slide_data["content"], list) else slide_data["content"]
                    slide.placeholders[1].text = content_text
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"{safe_title.replace(' ', '_')}_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        
        prs.save(filepath)
        
        print(f"✅ Custom presentation created!")
        print(f"📁 File: {filepath}")
        
        if collaborators:
            self._create_pptx_collaboration_guide(filepath, collaborators)
        
        return filepath
    
    def _create_pptx_collaboration_guide(self, pptx_file, collaborators):
        """Create collaboration guide for PowerPoint sharing"""
        
        base_name = os.path.splitext(os.path.basename(pptx_file))[0]
        guide_file = os.path.join(self.output_dir, f"{base_name}_COLLABORATION_GUIDE.md")
        
        with open(guide_file, 'w', encoding='utf-8') as f:
            f.write(f"""# PowerPoint Collaboration Guide
            
## 📋 Presentation File
`{os.path.basename(pptx_file)}`

## 🚀 Sharing Options

### Option 1: Email Attachment
1. **Attach .pptx file** to email
2. **Send to collaborators**: {', '.join(collaborators)}
3. **Recipients can open** in PowerPoint, Google Slides, or Keynote
4. **Edit directly** and send back changes

### Option 2: Cloud Collaboration (Recommended)
1. **Upload to OneDrive/Google Drive**
   - OneDrive: Real-time co-editing in PowerPoint Online
   - Google Drive: Opens in Google Slides with full editing
2. **Share link** with edit permissions
3. **Collaborate simultaneously** online

### Option 3: Microsoft Teams/SharePoint
1. **Upload to Teams** or SharePoint
2. **Collaborate in real-time** using PowerPoint Online
3. **Track changes** and version history
4. **Comment and review** features available

## 👥 Collaboration Workflow

### Immediate Actions:
- [ ] Upload file to preferred cloud platform
- [ ] Share with: {', '.join(collaborators)}
- [ ] Set edit permissions for all collaborators
- [ ] Send notification email with access instructions

### For Collaborators:
1. **Open shared file** via cloud link
2. **Edit directly** in browser or desktop app
3. **Add comments** for discussion points
4. **Use suggested edits** feature for tracked changes
5. **Notify team** when your section is complete

## 🛠️ Technical Details

### File Compatibility:
- ✅ **Microsoft PowerPoint** (2016 or later)
- ✅ **PowerPoint Online** (browser-based)
- ✅ **Google Slides** (automatic conversion)
- ✅ **Apple Keynote** (import with formatting preserved)
- ✅ **LibreOffice Impress** (open-source alternative)

### Features Available:
- **Real-time editing** (cloud platforms)
- **Version history** and restore
- **Comments and suggestions**
- **Presenter notes** editing
- **Slide transitions** and animations
- **Export to PDF** for final sharing

## 📧 Contact Information
**Questions or Issues:**
- Email: simonwang@hkbu.edu.hk
- Office: Language Centre, HKBU

## ✅ Next Steps
1. Choose collaboration platform (OneDrive recommended for PowerPoint)
2. Upload and share the presentation file
3. Send this guide to all collaborators
4. Set deadline for collaborative editing
5. Schedule review meeting for final approval
""")
        
        print(f"📋 Collaboration guide created: {os.path.basename(guide_file)}")

def main():
    parser = argparse.ArgumentParser(description="PowerPoint CLI Generator")
    parser.add_argument("--type", choices=["lang2077", "custom"], default="lang2077",
                       help="Type of presentation to create")
    parser.add_argument("--title", help="Custom presentation title")
    parser.add_argument("--collaborators", nargs="*", default=[],
                       help="Email addresses of collaborators")
    parser.add_argument("--open", action="store_true", help="Open generated presentation")
    
    args = parser.parse_args()
    
    ppt_cli = PowerPointCLI()
    
    if args.type == "lang2077":
        pptx_file = ppt_cli.create_lang2077_presentation(args.collaborators)
    else:
        if not args.title:
            args.title = input("Enter presentation title: ")
        # For demo, create a simple custom presentation
        slides_data = [
            {"type": "title", "title": args.title, "subtitle": "Created with PowerPoint CLI"},
            {"title": "Sample Content", "content": ["• Point 1", "• Point 2", "• Point 3"]}
        ]
        pptx_file = ppt_cli.create_custom_presentation(args.title, slides_data, args.collaborators)
    
    # Open presentation if requested
    if args.open:
        import subprocess
        subprocess.run(["open", pptx_file])  # macOS
        print("📊 Presentation opened")

if __name__ == "__main__":
    main()
